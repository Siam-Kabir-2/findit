"""
FindIt — clean university project presentation (PPTX).
White background, black text, dark-blue accent. Short text + screenshots.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
OUT = ROOT / "FindIt_Presentation.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

ACCENT = RGBColor(0x1E, 0x40, 0xAF)  # dark blue
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF1, 0xF5, 0xF9)
LINE = RGBColor(0xE2, 0xE8, 0xF0)


def run(p, text, size=18, bold=False, color=INK):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r


def bg(slide, color=BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    spTree = slide.shapes._spTree
    el = s._element
    spTree.remove(el)
    spTree.insert(2, el)


def accent_bar(slide):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.1), prs.slide_height)
    b.fill.solid()
    b.fill.fore_color.rgb = ACCENT
    b.line.fill.background()


def heading(slide, text, y=0.35):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.1), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    run(p, text, size=26, bold=True, color=INK)
    # underline accent
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y + 0.55), Inches(0.9), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def caption(slide, text, y=0.95):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.1), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    run(p, text, size=14, color=MUTED)


def bullets(slide, lines, top=1.3, left=0.6, width=12.0, size=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run(p, "•  " + line, size=size, color=MUTED)


def fit_picture(slide, name, left, top, max_w, max_h):
    path = SHOTS / name
    if not path.exists():
        return None
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top))
    ratio = min(Inches(max_w) / pic.width, Inches(max_h) / pic.height)
    pic.width = int(pic.width * ratio)
    pic.height = int(pic.height * ratio)
    return pic


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    accent_bar(s)
    return s


# ---------- slides ----------

def slide_title():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    t = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.2))
    run(t.text_frame.paragraphs[0], "FindIt", size=44, bold=True, color=INK)

    st = s.shapes.add_textbox(Inches(0.9), Inches(3.4), Inches(11.5), Inches(0.6))
    run(st.text_frame.paragraphs[0], "Lost and Found Item Management System", size=22, bold=True, color=ACCENT)

    sub = s.shapes.add_textbox(Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.8))
    tf = sub.text_frame
    tf.word_wrap = True
    run(tf.paragraphs[0], "A Web-Based Application Using Laravel and Oracle PL/SQL", size=16, color=MUTED)


def slide_overview():
    s = blank()
    heading(s, "Project Overview")
    bullets(s, [
        "Web-based lost and found system for campus users",
        "Users report items, search listings, and submit ownership claims",
        "Admins review claims and manage the full catalog",
        "Built with Laravel (web) and Oracle Database with PL/SQL (business logic)",
        "Replaces scattered social posts and notice boards with one platform",
    ])


def slide_problem():
    s = blank()
    heading(s, "Main Problem")
    bullets(s, [
        "Lost and found info is scattered across groups, boards, and staff",
        "Hard to search items and track what happened to them",
        "No clear way to verify ownership claims",
        "Little or no activity history",
    ], top=1.25)
    # solution card
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.6), Inches(12.0), Inches(2.0))
    card.fill.solid()
    card.fill.fore_color.rgb = SOFT
    card.line.color.rgb = LINE
    box = s.shapes.add_textbox(Inches(0.9), Inches(4.9), Inches(11.4), Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    run(tf.paragraphs[0], "FindIt solution", size=16, bold=True, color=ACCENT)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    run(p2, "One place to report, search, claim, and manage lost & found items.", size=16, color=MUTED)


def slide_objectives():
    s = blank()
    heading(s, "Project Objectives")
    # two columns
    left = s.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(5.8), Inches(5.5))
    tf = left.text_frame
    tf.word_wrap = True
    run(tf.paragraphs[0], "Users can", size=18, bold=True, color=ACCENT)
    for line in [
        "Register and log in",
        "Report lost / found items",
        "Search and filter items",
        "View item details",
        "Submit and track claims",
        "Manage their own items",
    ]:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        run(p, "•  " + line, size=15, color=MUTED)

    right = s.shapes.add_textbox(Inches(7.0), Inches(1.25), Inches(5.8), Inches(5.5))
    tf2 = right.text_frame
    tf2.word_wrap = True
    run(tf2.paragraphs[0], "Admins can", size=18, bold=True, color=ACCENT)
    for line in [
        "Review, approve, or reject claims",
        "Update item status",
        "Manage users",
        "Manage categories and locations",
        "View audit logs",
    ]:
        p = tf2.add_paragraph()
        p.space_before = Pt(6)
        run(p, "•  " + line, size=15, color=MUTED)


def slide_tech():
    s = blank()
    heading(s, "Technology Used")
    groups = [
        ("Frontend", "HTML · Blade · Tailwind CSS · JavaScript · Vite"),
        ("Backend", "PHP · Laravel 12"),
        ("Database", "Oracle 11g XE · SQL · PL/SQL (findit_pkg)"),
        ("Connection", "PDO_OCI · custom Laravel Oracle driver"),
        ("Tools", "VS Code · XAMPP · Composer · npm · SQL*Plus · GitHub"),
    ]
    y = 1.3
    for title, body in groups:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(12.0), Inches(0.95))
        card.fill.solid()
        card.fill.fore_color.rgb = SOFT
        card.line.color.rgb = LINE
        box = s.shapes.add_textbox(Inches(0.9), Inches(y + 0.18), Inches(11.4), Inches(0.7))
        tf = box.text_frame
        run(tf.paragraphs[0], title, size=15, bold=True, color=ACCENT)
        p = tf.add_paragraph()
        run(p, body, size=15, color=MUTED)
        y += 1.05


def slide_users():
    s = blank()
    heading(s, "System Users")
    cols = [
        ("Visitor", ["View and search items", "Open item details", "Register / log in"]),
        ("Registered User", ["Report lost / found items", "Upload images", "Submit claims", "Track my items & claims"]),
        ("Administrator", ["Separate admin login", "Approve / reject claims", "Manage catalog & users", "View audit logs"]),
    ]
    x = 0.6
    for title, items in cols:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(3.9), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = SOFT
        card.line.color.rgb = LINE
        box = s.shapes.add_textbox(Inches(x + 0.25), Inches(1.55), Inches(3.4), Inches(4.7))
        tf = box.text_frame
        tf.word_wrap = True
        run(tf.paragraphs[0], title, size=18, bold=True, color=ACCENT)
        for it in items:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
            run(p, "•  " + it, size=14, color=MUTED)
        x += 4.15


def slide_features_list():
    s = blank()
    heading(s, "Main Features")
    bullets(s, [
        "User registration and login",
        "Homepage and public item board",
        "Search and filter items",
        "Item details with images",
        "Report lost or found items",
        "User dashboard, My Items, My Claims",
        "Claim submission and status tracking",
        "Admin dashboard and claim management",
        "Item, user, category, and location management",
        "Oracle audit logs",
    ], size=16)


def slide_image(title, cap, image, max_h=5.5):
    s = blank()
    heading(s, title)
    caption(s, cap)
    fit_picture(s, image, 0.6, 1.4, 12.0, max_h)
    return s


def slide_two_images(title, cap, left_img, right_img, left_label="", right_label=""):
    s = blank()
    heading(s, title)
    caption(s, cap)
    if left_label:
        lb = s.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(5.8), Inches(0.3))
        run(lb.text_frame.paragraphs[0], left_label, size=13, bold=True, color=ACCENT)
    if right_label:
        rb = s.shapes.add_textbox(Inches(6.9), Inches(1.35), Inches(5.8), Inches(0.3))
        run(rb.text_frame.paragraphs[0], right_label, size=13, bold=True, color=ACCENT)
    y = 1.7 if (left_label or right_label) else 1.4
    fit_picture(s, left_img, 0.6, y, 5.8, 5.2)
    fit_picture(s, right_img, 6.9, y, 5.8, 5.2)
    return s


def slide_three_images(title, cap, imgs):
    s = blank()
    heading(s, title)
    caption(s, cap)
    x = 0.5
    for name in imgs:
        fit_picture(s, name, x, 1.45, 4.0, 5.3)
        x += 4.2
    return s


def slide_diagram(title, cap, image):
    return slide_image(title, cap, image, max_h=5.6)


def slide_architecture_text():
    s = blank()
    heading(s, "System Architecture")
    caption(s, "Request path from browser to Oracle.")
    fit_picture(s, "26-architecture.png", 0.6, 1.35, 12.0, 4.4)
    box = s.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12.0), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    run(tf.paragraphs[0], "Routes → middleware → controllers → models / PL/SQL service → PDO_OCI → findit_pkg → Oracle tables.", size=13, color=MUTED)


def slide_data_flows():
    s = blank()
    heading(s, "User and Admin Data Flow")
    # left card
    for x, title, lines in [
        (0.6, "User ↔ System", [
            "User sends: register, login, item data, search, claims",
            "System returns: results, details, dashboard, claim status",
        ]),
        (6.9, "Admin ↔ System", [
            "Admin sends: login, claim decisions, status updates, catalog changes",
            "System returns: stats, claims, items, users, audit records",
        ]),
    ]:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.4), Inches(5.8), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = SOFT
        card.line.color.rgb = LINE
        box = s.shapes.add_textbox(Inches(x + 0.3), Inches(1.7), Inches(5.2), Inches(4.2))
        tf = box.text_frame
        tf.word_wrap = True
        run(tf.paragraphs[0], title, size=18, bold=True, color=ACCENT)
        for line in lines:
            p = tf.add_paragraph()
            p.space_before = Pt(16)
            run(p, line, size=15, color=MUTED)


def slide_claim_approval():
    s = blank()
    heading(s, "Claim Approval Flow")
    bullets(s, [
        "Admin selects a pending claim",
        "System verifies the claim through PL/SQL",
        "Selected claim becomes Approved",
        "Item status becomes Claimed",
        "Other pending claims for that item become Rejected",
        "Audit log is written by Oracle triggers",
        "User and admin see the updated result",
    ])


def slide_auth_flow():
    s = blank()
    heading(s, "Authentication Flow")
    bullets(s, [
        "User enters email and password",
        "Laravel validates the input",
        "System looks up the account in Oracle",
        "Password is verified (hashed)",
        "Session is created",
        "User dashboard opens",
        "Admins use the same idea with a separate login and admin guard",
    ])


def slide_db_overview():
    s = blank()
    heading(s, "Database Overview")
    caption(s, "Seven main tables store all FindIt data.")
    fit_picture(s, "33-oracle-tables.png", 0.6, 1.35, 12.0, 5.5)


def slide_tables():
    s = blank()
    heading(s, "Database Tables")
    bullets(s, [
        "USERS — registered campus users",
        "ADMINS — administrator accounts",
        "CATEGORIES — item categories",
        "LOCATIONS — campus places",
        "ITEMS — lost / found reports (type, status, image)",
        "CLAIMS — ownership requests (pending / approved / rejected)",
        "AUDIT_LOGS — trigger-written activity history",
    ], size=16)


def slide_relationships():
    s = blank()
    heading(s, "Database Relationships")
    bullets(s, [
        "One user can report many items",
        "Each item belongs to one category and one location",
        "One user can submit many claims",
        "Each claim belongs to one user and one item",
        "One item can receive many claims",
        "Audit logs record changes on items and claims",
    ])


def slide_plsql():
    s = blank()
    heading(s, "PL/SQL Package — FINDIT_PKG")
    caption(s, "Important business rules live in Oracle.")
    fit_picture(s, "34-plsql-package.png", 0.5, 1.3, 12.2, 5.6)


def slide_sequences():
    s = blank()
    heading(s, "Sequences and Triggers")
    bullets(s, [
        "Sequences generate unique IDs for every table",
        "Insert triggers assign the next sequence value",
        "Audit triggers log insert / update / delete on items and claims",
    ], top=1.25)
    # mini flow cards
    steps = ["New record", "Trigger runs", "Sequence ID", "Record saved"]
    x = 0.6
    for i, step in enumerate(steps):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.2), Inches(2.6), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = SOFT
        card.line.color.rgb = LINE
        box = s.shapes.add_textbox(Inches(x + 0.15), Inches(4.65), Inches(2.3), Inches(0.7))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, step, size=15, bold=True, color=ACCENT)
        if i < len(steps) - 1:
            arrow = s.shapes.add_textbox(Inches(x + 2.55), Inches(4.65), Inches(0.4), Inches(0.5))
            ap = arrow.text_frame.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            run(ap, "→", size=20, bold=True, color=ACCENT)
        x += 3.1


def slide_controllers():
    s = blank()
    heading(s, "Controllers and Models")
    left = s.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(6.0), Inches(5.5))
    tf = left.text_frame
    tf.word_wrap = True
    run(tf.paragraphs[0], "Controllers", size=17, bold=True, color=ACCENT)
    for line in [
        "User: Home, Item, Dashboard, Claim, Auth",
        "Admin: Auth, Dashboard, Claim, Item,",
        "User, Category, Location, Audit",
        "",
        "Controllers validate input, call services,",
        "and return Blade views.",
    ]:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        run(p, line, size=14, color=MUTED)

    right = s.shapes.add_textbox(Inches(7.0), Inches(1.25), Inches(5.8), Inches(5.5))
    tf2 = right.text_frame
    tf2.word_wrap = True
    run(tf2.paragraphs[0], "Models", size=17, bold=True, color=ACCENT)
    for line in ["User", "Admin", "Item", "Claim", "Category", "Location", "AuditLog"]:
        p = tf2.add_paragraph()
        p.space_before = Pt(6)
        run(p, "•  " + line, size=15, color=MUTED)


def slide_security():
    s = blank()
    heading(s, "Security and Validation")
    bullets(s, [
        "Password hashing and session-based auth",
        "Separate user and admin guards / routes",
        "Form validation and image file checks",
        "Oracle constraints, foreign keys, unique email",
        "Allowed item types and status values only",
        "Parameterized queries through PDO / PL/SQL",
        "Users cannot claim their own item or duplicate pending claims",
    ])


def slide_setup():
    s = blank()
    heading(s, "Project Setup")
    bullets(s, [
        "Install PHP, Composer, Node.js, Oracle XE, Instant Client",
        "Clone the GitHub repository",
        "composer install  →  copy .env  →  configure Oracle",
        "Run database/oracle SQL scripts in order",
        "npm install && npm run build",
        "php artisan storage:link",
        "php artisan serve  →  open in browser",
    ], size=16)


def slide_oracle_files():
    s = blank()
    heading(s, "Oracle SQL Files")
    caption(s, "Scripts that prepare the full database.")
    fit_picture(s, "32-oracle-sql-files.png", 0.6, 1.3, 12.0, 5.6)


def slide_testing():
    s = blank()
    heading(s, "Testing")
    caption(s, "Sample test cases used during demonstration.")
    # simple table-like bullets in two columns
    left_tests = [
        "Register / login / invalid login",
        "Submit lost and found items",
        "Search and filter items",
        "View details and submit claim",
        "Block own-item and duplicate claims",
    ]
    right_tests = [
        "Approve / reject claim",
        "Update item status",
        "Add category / location",
        "View audit logs",
        "Block unauthorized admin access",
    ]
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.0))
    tf = box.text_frame
    for i, t in enumerate(left_tests):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run(p, "•  " + t + "  →  Passed", size=14, color=MUTED)
    box2 = s.shapes.add_textbox(Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.0))
    tf2 = box2.text_frame
    for i, t in enumerate(right_tests):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(8)
        run(p, "•  " + t + "  →  Passed", size=14, color=MUTED)


def slide_benefits():
    s = blank()
    heading(s, "Benefits")
    bullets(s, [
        "Centralized lost and found information",
        "Easy search and filtering",
        "Structured claims with clear statuses",
        "Separate user and admin access",
        "Oracle integrity plus audit history",
        "Less manual / informal follow-up work",
    ])


def slide_result():
    s = blank()
    heading(s, "Project Result")
    bullets(s, [
        "Working campus lost and found web platform",
        "Users can report, search, claim, and track decisions",
        "Admins manage the full process from a separate console",
        "Laravel handles the UI and requests",
        "Oracle PL/SQL handles core database operations",
    ])


def slide_conclusion():
    s = blank()
    heading(s, "Conclusion")
    bullets(s, [
        "FindIt organizes the full lost-and-found workflow",
        "From reporting an item to approving a claim",
        "Shows Laravel web development with Oracle PL/SQL",
        "Includes auth, roles, schema design, triggers, and audit logs",
    ])


def slide_thanks():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0))
    run(t.text_frame.paragraphs[0], "Thank You", size=44, bold=True, color=INK)
    st = s.shapes.add_textbox(Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.5))
    run(st.text_frame.paragraphs[0], "Questions and Discussion", size=20, color=ACCENT)
    sub = s.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.8))
    run(sub.text_frame.paragraphs[0], "FindIt – Lost and Found Item Management System", size=15, color=MUTED)


# ===== build order (~35–38 slides) =====
slide_title()                          # 1
slide_overview()                       # 2
slide_problem()                        # 3
slide_objectives()                     # 4
slide_tech()                           # 5
slide_users()                          # 6
slide_diagram("System Workflow", "End-to-end process from visit to claim result.", "29-workflow.png")  # 7
slide_architecture_text()              # 8

slide_image("Homepage", "Landing page — browse, login, and register.", "20-home-viewport.png")  # 9
slide_two_images(
    "Browse and Search",
    "Public board with lost / found filters.",
    "19-browse-viewport.png",
    "21-browse-lost.png",
    "All items",
    "Filtered: Lost",
)  # 10
slide_image("Item Details", "Name, description, image, location, status, and claim option.", "03-item-detail.png")  # 11
slide_two_images(
    "Registration and Login",
    "Create an account, then sign in to the user area.",
    "05-register.png",
    "04-user-login.png",
    "Register",
    "Login",
)  # 12
slide_image("User Dashboard", "Summary of the user’s items and claims.", "06-user-dashboard.png")  # 13
slide_two_images(
    "Report an Item",
    "Submission form and the user’s item list after reporting.",
    "07-report-item.png",
    "08-my-items.png",
    "Add item form",
    "My Items",
)  # 14
slide_two_images(
    "Submit a Claim",
    "Claim form with message and proof, then My Claims list.",
    "23-claim-form.png",
    "09-my-claims.png",
    "Claim form",
    "My Claims",
)  # 15

slide_two_images(
    "Admin Login and Dashboard",
    "Separate admin login and live system statistics.",
    "10-admin-login.png",
    "11-admin-dashboard.png",
    "Admin login",
    "Admin dashboard",
)  # 16
slide_image("Claim Management", "Admins review pending claims and approve or reject them.", "25-admin-claims-pending.png")  # 17
slide_image("Item Management", "View items and update status from the admin console.", "14-admin-items.png")  # 18
slide_image("User Management", "View and manage registered users.", "15-admin-users.png")  # 19
slide_two_images(
    "Categories and Locations",
    "Catalog data used when reporting items.",
    "16-admin-categories.png",
    "17-admin-locations.png",
    "Categories",
    "Locations",
)  # 20
slide_image("Audit Logs", "Trigger-written history of item and claim changes.", "18-admin-audit.png")  # 21

slide_data_flows()                     # 22
slide_claim_approval()                 # 23
slide_auth_flow()                      # 24
slide_db_overview()                    # 25
slide_diagram(
    "Entity Relationship Diagram",
    "USERS, ITEMS, CLAIMS, CATEGORIES, LOCATIONS, ADMINS, AUDIT_LOGS.",
    "27-er-diagram.png",
)  # 26
slide_diagram(
    "Database Schema Diagram",
    "Primary keys, foreign keys, and table links.",
    "28-schema-diagram.png",
)  # 27
slide_tables()                         # 28
slide_relationships()                  # 29
slide_plsql()                          # 30
slide_sequences()                      # 31
slide_controllers()                    # 32
slide_security()                       # 33
slide_setup()                          # 34
slide_oracle_files()                   # 35
slide_image("Project Folder Structure", "Main Laravel folders used in the project.", "31-folder-structure.png")  # 36
slide_image("GitHub Repository", "Project source on GitHub.", "30-github-repo.png")  # 37
slide_testing()                        # 38
slide_benefits()                       # 39
slide_conclusion()                     # 40
slide_thanks()                         # 41

# Combined: dropped redundant features list + duplicate my-items/my-claims/result slides.

try:
    prs.save(OUT)
except PermissionError:
    alt = ROOT / "FindIt_Presentation_v2.pptx"
    prs.save(alt)
    print(f"Original locked. Saved: {alt}")
    print(f"Slides: {len(prs.slides)}")
else:
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
